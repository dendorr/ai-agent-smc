"""
DOCUMENTS AGENT v5 — Universal document intelligence (async + GLM-OCR vision)

Tutte le funzioni di estrazione (PPTX/DOCX/PDF) restano sincrone perché
non beneficerebbero dall'async (sono CPU-bound + I/O su file locali).
Le chiamate LLM e ChromaDB sono async.

Supporta: PDF (anche scansionati), PPTX (con OCR immagini), DOCX, Markdown, TXT.

Architettura:
  - Markdown cache: convert one-time, reuse on restart (mtime-based)
  - Image OCR a 3 livelli con fallback automatico:
      Livello 0: GLM-OCR (vision multimodale via Ollama) — top accuracy
      Livello 1: pytesseract — veloce, CPU, fallback locale
      Livello 2: easyocr     — neurale, fallback finale
  - PDF scansionati: detection automatica + rasterizzazione pagina intera
                     + GLM-OCR (estrae testo, tabelle e formule come Markdown)
  - Multi-model: routing model (LLM_MODEL_FAST) seleziona documenti,
                 answer model (LLM_MODEL_MAIN) genera la risposta
  - Persistent memory: indice documenti + annotazioni utente
  - Zero hardcoded: nessuna assunzione su contenuto/lingua/dominio
  - Air-gapped: GLM-OCR gira via Ollama in locale, niente cloud
"""

import sys
import os
import json
import re
import hashlib
import asyncio
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config.config import (
    CHROMA_PATHS, LLM_MODEL_MAIN, LLM_MODEL_FAST,
    CHUNK_SIZE, CHUNK_OVERLAP, EXTENSIONS,
    MEMORY_PATH, MARKDOWN_CACHE_DIR,
    OCR_ENABLED, OCR_MODEL, OCR_MIN_TEXT_LEN,
    OCR_PDF_PAGE_RASTER_DPI, OCR_PDF_MIN_TEXT_PER_PAGE,
)

import chromadb
import semantic_analyzer as analyzer
from llm_client import chat_complete, chat_complete_json, vision_extract_sync

client     = chromadb.PersistentClient(path=CHROMA_PATHS["documents"])
collection = client.get_or_create_collection("documents")

MEMORY_FILE    = Path(MEMORY_PATH) / "documents_memory.json"
MARKDOWN_CACHE = Path(MARKDOWN_CACHE_DIR)

# ── Models ────────────────────────────────────────────────────────────────────
ROUTING_MODEL = LLM_MODEL_FAST   # fast: seleziona documenti rilevanti
ANSWER_MODEL  = LLM_MODEL_MAIN   # accurate: risposta finale

# ── OCR — singleton (no re-init) ──────────────────────────────────────────────
_easyocr_reader = None


def _get_easyocr_reader():
    """Lazy-init del Reader easyocr — una sola volta per processo."""
    global _easyocr_reader
    if _easyocr_reader is None:
        try:
            import easyocr
            _easyocr_reader = easyocr.Reader(["it", "en"], gpu=False, verbose=False)
        except Exception:
            pass
    return _easyocr_reader


# ── GLM-OCR vision (Livello 0) ────────────────────────────────────────────────

def ocr_with_glm(image_bytes: bytes, source_hint: str = "") -> str:
    """
    OCR via GLM-OCR multimodale (Ollama, locale).

    Si attende image_bytes in formato PNG (canonico). Ritorna il testo
    estratto pulito o stringa vuota in caso di errore / OCR disabilitato.
    Non solleva eccezioni: il chiamante può fare fallback su altri livelli.
    """
    if not image_bytes or not OCR_ENABLED:
        return ""
    try:
        text = vision_extract_sync(
            model=OCR_MODEL,
            image_bytes=image_bytes,
            prompt="Text Recognition:",
            image_format="png",
        )
        return text.strip() if text else ""
    except Exception:
        return ""


def ocr_image_bytes(image_bytes: bytes, source_hint: str = "") -> str:
    """
    OCR su bytes raw, con cascata a 3 livelli e fallback automatico.

      Livello 0: GLM-OCR (vision multimodale via Ollama) — top accuracy
                 su layout complessi, tabelle, formule. Skippato se OCR_ENABLED=False.
      Livello 1: pytesseract — veloce, CPU, accurato per testo stampato/scannerizzato
      Livello 2: easyocr     — neurale, migliore per testo stilizzato/manoscritto

    Restituisce il blocco di testo estratto, o source_hint se nulla è leggibile.
    """
    if not image_bytes:
        return ""

    from PIL import Image
    import io

    try:
        img = Image.open(io.BytesIO(image_bytes))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
    except Exception:
        return ""

    # Conversione canonica a PNG bytes — formato consistente per tutti i livelli
    try:
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        png_bytes = buf.getvalue()
    except Exception:
        png_bytes = image_bytes  # fallback ai bytes originali

    # Livello 0: GLM-OCR
    if OCR_ENABLED:
        try:
            text = ocr_with_glm(png_bytes, source_hint=source_hint)
            if text and len(text) > OCR_MIN_TEXT_LEN:
                return f"[OCR]\n{text}"
        except Exception:
            pass

    # Livello 1: pytesseract
    try:
        import pytesseract
        text = pytesseract.image_to_string(
            img, lang="ita+eng", config="--psm 3 --oem 3"
        ).strip()
        if len(text) > OCR_MIN_TEXT_LEN:
            return f"[OCR]\n{text}"
    except Exception:
        pass

    # Livello 2: easyocr
    try:
        import numpy as np
        reader = _get_easyocr_reader()
        if reader:
            results = reader.readtext(np.array(img))
            lines = [r[1] for r in results if len(r) > 2 and r[2] > 0.3]
            text = "\n".join(lines).strip()
            if len(text) > OCR_MIN_TEXT_LEN:
                return f"[OCR]\n{text}"
    except Exception:
        pass

    return source_hint


def ocr_image_file(filepath) -> str:
    """OCR su un percorso file."""
    try:
        return ocr_image_bytes(Path(filepath).read_bytes())
    except Exception:
        return ""


# ── Helper condivisi ──────────────────────────────────────────────────────────

def _table_to_markdown(table) -> str:
    """
    Converte un Table object (python-pptx o python-docx) in tabella Markdown.
    Preserva tutto il contenuto cella esattamente — niente arrotondamenti.
    """
    rows = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            rows.append("|" + "|".join(["---"] * len(cells)) + "|")
    return "\n".join(rows)


# ── PPTX → Markdown ───────────────────────────────────────────────────────────

def convert_pptx_to_markdown(filepath) -> str:
    """
    Converte un PPTX in Markdown strutturato.

    Per slide:
      - Numero + titolo come heading
      - Text box (indentazione level-aware)
      - Tabelle → Markdown (numeri/unità esatti)
      - Immagini embedded → OCR
      - Note relatore → blockquote a fine slide
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    p   = Path(filepath)
    prs = Presentation(str(filepath))

    lines = [f"# {p.stem}", f"*File: {p.name} — {len(prs.slides)} slide*", ""]

    for slide_num, slide in enumerate(prs.slides, 1):
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except Exception:
            pass

        slide_heading = f"## Slide {slide_num}"
        if title:
            slide_heading += f": {title}"
        lines.append(slide_heading)
        lines.append("")

        for shape in slide.shapes:
            if shape.has_table:
                lines.append(_table_to_markdown(shape.table))
                lines.append("")
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text or text == title:
                        continue
                    level  = getattr(para, "level", 0)
                    indent = "  " * level
                    prefix = f"{'#' * (level + 3)} " if level > 0 else ""
                    lines.append(f"{indent}{prefix}{text}")
                lines.append("")
                continue

            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_bytes = shape.image.blob
                    hint = f"[Immagine — slide {slide_num}]"
                    ocr_text = ocr_image_bytes(img_bytes, source_hint=hint)
                    if ocr_text:
                        quoted = ocr_text.replace("\n", "\n> ")
                        lines.append(f"> **{hint}**")
                        lines.append(f"> {quoted}")
                        lines.append("")
            except Exception:
                pass

        try:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip()
            if notes_text:
                quoted = notes_text.replace("\n", "\n> ")
                lines.append("**Note relatore:**")
                lines.append(f"> {quoted}")
                lines.append("")
        except Exception:
            pass

        lines.append("---")
        lines.append("")

    return "\n".join(lines)


# ── DOCX → Markdown ───────────────────────────────────────────────────────────

def convert_docx_to_markdown(filepath) -> str:
    """
    Converte un DOCX in Markdown strutturato.

    Strategia:
      Livello 1 — markitdown (Microsoft): heading/tabelle/liste perfetti
      Livello 2 — python-docx manuale: itera body in ordine documento
                  (paragrafi + tabelle interleaved), poi estrae OCR immagini

    Numeri, unità, tolleranze sempre preservati esattamente.
    """
    p = Path(filepath)

    # Livello 1: markitdown
    try:
        from markitdown import MarkItDown
        result = MarkItDown().convert(str(filepath))
        text = result.text_content.strip()
        if len(text) > 100:
            md = f"# {p.stem}\n*File: {p.name}*\n\n{text}"
            img_blocks = _docx_extract_image_ocr(filepath)
            if img_blocks:
                md += "\n\n## Immagini estratte dal documento\n\n" + "\n\n".join(img_blocks)
            return md
    except ImportError:
        pass
    except Exception as e:
        print(f"  [markitdown warning] {p.name}: {e}", flush=True)

    # Livello 2: python-docx manuale
    try:
        from docx import Document
        from docx.text.paragraph import Paragraph
        from docx.table import Table as DocxTable

        doc = Document(str(filepath))
        lines = [f"# {p.stem}", f"*File: {p.name}*", ""]

        _HEADING_MAP = {
            "heading 1": "#", "heading 2": "##", "heading 3": "###",
            "heading 4": "####", "heading 5": "#####",
        }

        for child in doc.element.body.iterchildren():
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "p":
                para = Paragraph(child, doc)
                text = para.text.strip()
                style = (para.style.name or "").lower() if para.style else ""

                if not text:
                    lines.append("")
                    continue

                prefix = ""
                for key, md_prefix in _HEADING_MAP.items():
                    if key in style:
                        prefix = md_prefix
                        break

                if prefix:
                    lines.append(f"{prefix} {text}")
                elif "list" in style:
                    indent = "  " * (int(re.search(r'\d', style).group()) - 1
                                     if re.search(r'\d', style) else 0)
                    lines.append(f"{indent}- {text}")
                else:
                    lines.append(text)

            elif tag == "tbl":
                tbl = DocxTable(child, doc)
                lines.append("")
                lines.append(_table_to_markdown(tbl))
                lines.append("")

        lines.append("")

        img_blocks = _docx_extract_image_ocr(filepath)
        if img_blocks:
            lines.append("## Immagini estratte dal documento")
            lines.extend(img_blocks)

        return "\n".join(lines)

    except Exception as e:
        return f"[DOCX error: {e}]"


def _docx_extract_image_ocr(filepath) -> list:
    """Estrae tutte le immagini embedded nel DOCX e fa OCR su ognuna."""
    blocks = []
    try:
        from docx import Document
        doc = Document(str(filepath))
        for i, rel in enumerate(doc.part.rels.values(), 1):
            if "image" in rel.reltype:
                try:
                    img_bytes = rel.target_part.blob
                    ocr_text = ocr_image_bytes(
                        img_bytes, source_hint=f"[Immagine {i}]"
                    )
                    if ocr_text and len(ocr_text) > 10:
                        quoted = ocr_text.replace("\n", "\n> ")
                        blocks.append(f"> **Immagine {i}:**\n> {quoted}")
                except Exception:
                    pass
    except Exception:
        pass
    return blocks


# ── PDF → Markdown ────────────────────────────────────────────────────────────

def convert_pdf_to_markdown(filepath) -> str:
    """
    Converte un PDF in Markdown strutturato.

    Livello 1 — markitdown:  miglior preservazione layout su PDF nativi
    Livello 2 — fitz:        testo + OCR pagine scansionate (GLM-OCR)
                             + OCR immagini embedded + tabelle pagina-per-pagina
    Livello 3 — pdfplumber:  tabelle supplementari

    PDF scansionati: pagine con < OCR_PDF_MIN_TEXT_PER_PAGE caratteri vengono
    rasterizzate a OCR_PDF_PAGE_RASTER_DPI DPI e mandate intere a GLM-OCR.

    Tutti i valori numerici, unità e tolleranze preservati esattamente.
    """
    import fitz as _fitz

    p = Path(filepath)
    header = f"# {p.stem}\n*File: {p.name}*\n\n"

    # Livello 1: markitdown
    try:
        from markitdown import MarkItDown
        result = MarkItDown().convert(str(filepath))
        text = result.text_content.strip()
        if len(text) > 100:
            return header + text
    except ImportError:
        pass
    except Exception as e:
        print(f"  [markitdown warning] {p.name}: {e}", flush=True)

    # Livello 2: fitz
    fitz_parts = []
    try:
        doc = _fitz.open(str(filepath))
        for page_num, page in enumerate(doc, 1):
            page_lines = [f"## Pagina {page_num}", ""]

            text = page.get_text("text")
            text_clean = text.strip()
            page_was_ocr = False  # True = pagina gestita come scansione (GLM-OCR)

            # Pagina con testo nativo: usa quello
            if len(text_clean) >= OCR_PDF_MIN_TEXT_PER_PAGE:
                page_lines.append(text_clean)
                page_lines.append("")

            # Pagina con poco testo + OCR abilitato: probabile scansione
            elif OCR_ENABLED:
                try:
                    # Rasterizzazione pagina intera (PDF base = 72 DPI)
                    zoom = OCR_PDF_PAGE_RASTER_DPI / 72
                    mat  = _fitz.Matrix(zoom, zoom)
                    pix  = page.get_pixmap(matrix=mat)
                    page_png = pix.tobytes("png")

                    ocr_text = ocr_with_glm(
                        page_png,
                        source_hint=f"[Pagina {page_num} (scansione)]",
                    )
                    if ocr_text and len(ocr_text) > OCR_MIN_TEXT_LEN:
                        page_lines.append(f"[OCR pagina {page_num} — scansione]")
                        page_lines.append(ocr_text)
                        page_lines.append("")
                        page_was_ocr = True
                    elif text_clean:
                        # GLM-OCR fallito: usa il poco testo nativo come fallback
                        page_lines.append(text_clean)
                        page_lines.append("")
                except Exception as e:
                    print(f"  [GLM-OCR warning] pag. {page_num}: {e}", flush=True)
                    if text_clean:
                        page_lines.append(text_clean)
                        page_lines.append("")

            # OCR disabilitato: usa il poco testo nativo se presente
            elif text_clean:
                page_lines.append(text_clean)
                page_lines.append("")

            # Tabelle via fitz (PyMuPDF >= 1.23)
            try:
                for tab in page.find_tables().tables:
                    df = tab.to_pandas()
                    md = df.to_markdown(index=False)
                    if md:
                        page_lines.append(f"**Tabella:**\n{md}")
                        page_lines.append("")
            except Exception:
                pass

            # Immagini embedded → OCR
            # Skip se la pagina è già stata gestita via rasterizzazione GLM-OCR:
            # le immagini embedded sono già comprese nel rasterizzato — evita
            # duplicati in ChromaDB.
            if not page_was_ocr:
                for img_info in page.get_images(full=True):
                    try:
                        xref = img_info[0]
                        base_img = doc.extract_image(xref)
                        img_bytes = base_img.get("image", b"")
                        ocr_text = ocr_image_bytes(
                            img_bytes, source_hint=f"[Immagine pag. {page_num}]"
                        )
                        if ocr_text and len(ocr_text) > OCR_MIN_TEXT_LEN:
                            quoted = ocr_text.replace("\n", "\n> ")
                            page_lines.append(f"> **Immagine pag. {page_num}:**\n> {quoted}")
                            page_lines.append("")
                    except Exception:
                        pass

            fitz_parts.append("\n".join(page_lines))
        doc.close()
    except Exception as e:
        fitz_parts.append(f"[fitz error: {e}]")

    # Livello 3: pdfplumber
    plumber_parts = []
    try:
        import pdfplumber
        with pdfplumber.open(str(filepath)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                for table in page.extract_tables() or []:
                    rows = []
                    for i, row in enumerate(table):
                        cells = [str(c).strip() if c else "" for c in row]
                        rows.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            rows.append("|" + "|".join(["---"] * len(cells)) + "|")
                    if rows:
                        plumber_parts.append(
                            f"**Tabella pag. {page_num}:**\n" + "\n".join(rows)
                        )
    except Exception:
        pass

    result = header
    if fitz_parts:
        result += "\n\n---\n\n".join(fitz_parts)
    if plumber_parts:
        result += "\n\n## Tabelle aggiuntive\n\n" + "\n\n".join(plumber_parts)

    return result if len(result) > len(header) + 20 else f"[Impossibile estrarre {p.name}]"


# ── Markdown cache ────────────────────────────────────────────────────────────

def _cache_path(filepath) -> Path:
    """Path deterministico nella cache per un file sorgente."""
    h = hashlib.md5(str(filepath).encode()).hexdigest()[:8]
    stem = re.sub(r"[^a-zA-Z0-9]", "_", Path(filepath).stem[:40]).strip("_")
    return MARKDOWN_CACHE / f"{stem}_{h}.md"


def _cache_is_valid(filepath, cache_file: Path) -> bool:
    """Cache hit: file esiste E è più recente (o uguale) del sorgente."""
    if not cache_file.exists():
        return False
    try:
        return cache_file.stat().st_mtime >= Path(filepath).stat().st_mtime
    except Exception:
        return False


def get_or_create_markdown(filepath) -> str:
    """
    Restituisce il Markdown del documento.
    Legge dalla cache se valida; altrimenti converte e salva la cache.
    Invalidazione automatica: rigenera quando il file sorgente cambia.
    """
    MARKDOWN_CACHE.mkdir(parents=True, exist_ok=True)
    cache_file = _cache_path(filepath)

    if _cache_is_valid(filepath, cache_file):
        try:
            return cache_file.read_text(encoding="utf-8")
        except Exception:
            pass

    ext = Path(filepath).suffix.lower()
    markdown = ""

    if ext == ".pdf":
        markdown = convert_pdf_to_markdown(filepath)
    elif ext in (".pptx", ".ppt"):
        markdown = convert_pptx_to_markdown(filepath)
    elif ext in (".docx", ".doc"):
        markdown = convert_docx_to_markdown(filepath)
    elif ext == ".md":
        markdown = Path(filepath).read_text(encoding="utf-8", errors="ignore")
    elif ext == ".txt":
        content = Path(filepath).read_text(encoding="utf-8", errors="ignore")
        markdown = f"# {Path(filepath).stem}\n\n{content}"

    if markdown and len(markdown.strip()) > 50:
        try:
            cache_file.write_text(markdown, encoding="utf-8")
            print(
                f"  [cache] Scritto {cache_file.name} ({len(markdown):,} caratteri)",
                flush=True,
            )
        except Exception as e:
            print(f"  [cache warning] {e}", flush=True)

    return markdown


# ── Memoria documenti ─────────────────────────────────────────────────────────

def load_memory() -> dict:
    if MEMORY_FILE.exists():
        try:
            return json.loads(MEMORY_FILE.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {"documents": {}, "annotations": {}}


def save_memory(memory: dict):
    MEMORY_FILE.parent.mkdir(parents=True, exist_ok=True)
    MEMORY_FILE.write_text(
        json.dumps(memory, indent=2, ensure_ascii=False), encoding="utf-8"
    )


def update_document_memory(filepath, markdown: str):
    """Salva metadata per-documento dopo l'indicizzazione."""
    from datetime import datetime
    memory = load_memory()
    if "documents" not in memory:
        memory["documents"] = {}

    p = Path(filepath)
    memory["documents"][p.name] = {
        "filepath":   str(filepath),
        "indexed_at": datetime.now().isoformat(timespec="seconds"),
        "words":      len(markdown.split()),
        "chars":      len(markdown),
        "size_kb":    round(p.stat().st_size / 1024, 1),
        "ext":        p.suffix.lower(),
    }
    save_memory(memory)


def add_annotation(filename: str, note: str):
    """Allega un'annotazione utente a un documento (persistente)."""
    memory = load_memory()
    if "annotations" not in memory:
        memory["annotations"] = {}
    if filename not in memory["annotations"]:
        memory["annotations"][filename] = []
    memory["annotations"][filename].append(note)
    save_memory(memory)


# ── Chunking + indexing (sync — chiamato dal watcher sync) ────────────────────

def chunk_text(text: str) -> list:
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunks.append(" ".join(words[i:i + CHUNK_SIZE]))
        i += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks or [""]


def index_file(filepath) -> int:
    """
    Converte un documento in Markdown (o usa cache), spezza in chunk e
    salva su ChromaDB. Ritorna il numero di chunk indicizzati, 0 se saltato.
    Funzione sincrona — il watcher è sync e gira fuori dall'event loop.
    """
    p = Path(filepath)
    ext = p.suffix.lower()

    if ext not in EXTENSIONS["documents"]:
        return 0

    print(f"  [docs] Elaborazione {p.name}...", flush=True)
    markdown = get_or_create_markdown(filepath)

    if not markdown or not markdown.strip():
        return 0

    update_document_memory(filepath, markdown)

    # Rimuovi chunk stale
    try:
        existing = collection.get(where={"path": str(filepath)})
        if existing and existing["ids"]:
            collection.delete(ids=existing["ids"])
    except Exception:
        pass

    chunks = chunk_text(markdown)
    for i, chunk in enumerate(chunks):
        collection.upsert(
            documents=[chunk],
            ids=[f"{filepath}__c{i}"],
            metadatas=[{
                "filename": p.name,
                "path":     str(filepath),
                "chunk":    i,
                "agent":    "documents",
                "type":     "chunk",
                "ext":      ext,
            }]
        )
    return len(chunks)


def index_folder(folder):
    files = [f for f in Path(folder).rglob("*")
             if f.is_file() and f.suffix.lower() in EXTENSIONS["documents"]]
    print(f"[Documents] Trovati {len(files)} file...")
    total = 0
    for f in files:
        n = index_file(str(f))
        if n: print(f"  [OK] {f.name} → {n} chunk")
        else: print(f"  [--] {f.name} → saltato")
        total += n
    print(f"[Documents] Completato — {total} chunk totali")

# ── Filename detection ────────────────────────────────────────────────────────

def detect_filename_filter(query: str) -> dict | None:
    """
    Analizza la query dell'utente per individuare riferimenti a file specifici.
    Se trova un match, ritorna un filtro ChromaDB 'where' per filename.
    Altrimenti ritorna None (ricerca semantica normale).
    
    Gestisce varianti come: "lez 13", "Lez13", "lezione 13", "file lez13",
    "prova itinere 1", "ProvaItinere1", "mock exam", ecc.
    """
    import chromadb

    # Prendi tutti i filename indicizzati
    try:
        all_meta = collection.get(include=["metadatas"])
        all_filenames = list({
            m["filename"] for m in all_meta["metadatas"]
            if m.get("filename") and m.get("type") != "semantic_card"
        })
    except Exception:
        return None

    if not all_filenames:
        return None

    # Normalizza: rimuovi estensione, spazi, underscore, tutto lowercase
    def normalize(s: str) -> str:
        s = Path(s).stem if "." in s else s
        # Separa camelCase/PascalCase: "ProvaItinere1" → "prova itinere 1"
        s = re.sub(r'([a-z])([A-Z])', r'\1 \2', s)
        # Separa lettere da numeri: "Lez13" → "Lez 13"
        s = re.sub(r'([a-zA-Z])(\d)', r'\1 \2', s)
        s = re.sub(r'(\d)([a-zA-Z])', r'\1 \2', s)
        return re.sub(r'[_\-\s]+', ' ', s).strip().lower()

    query_norm = normalize(query)

    # Cerca il filename con il miglior match
    best_match = None
    best_score = 0

    for fname in all_filenames:
        fname_norm = normalize(fname)

        # Match esatto della parte normalizzata
        if fname_norm in query_norm or query_norm in fname_norm:
            score = len(fname_norm)
            if score > best_score:
                best_score = score
                best_match = fname
            continue

        # Match parziale: tutte le parole del filename presenti nella query
        fname_words = fname_norm.split()
        query_words = query_norm.split()
        if len(fname_words) >= 1 and all(fw in query_words for fw in fname_words):
            score = len(fname_norm)
            if score > best_score:
                best_score = score
                best_match = fname

    if best_match:
        print(f"  [smart-filter] Query '{query}' → filtro per '{best_match}'", flush=True)
        return {"filename": best_match}

    return None

# ── Search (async — called from the server as async) ────────────────────────────────

async def search(query: str) -> str:
    """Async vector DB search with smart filename filtering + lazy semantic cards."""
    where_filter = detect_filename_filter(query)
    return await analyzer.search_with_cards(
        collection, query, "documents", n_results=6, where_filter=where_filter
    )


# ── Routing model (async) ─────────────────────────────────────────────────────

async def route_documents(query: str, memory: dict) -> list:
    """
    Usa il modello fast per identificare i documenti più probabilmente rilevanti.
    Ritorna lista di filename. In caso di errore, lista vuota.
    """
    docs = memory.get("documents", {})
    if not docs:
        return []

    doc_lines = "\n".join(
        f"  - {name}  ({info.get('words', 0):,} parole | "
        f"{info.get('size_kb', 0)} KB | {info.get('ext', '')})"
        for name, info in list(docs.items())[:25]
    )

    system = "Sei un assistente che seleziona documenti rilevanti. Rispondi solo con nomi file, uno per riga."
    user = f"""DOCUMENTI DISPONIBILI:
{doc_lines}

DOMANDA: {query}

Elenca SOLO i nomi file più rilevanti (massimo 4), uno per riga.
Solo nomi file esatti, nessun altro testo."""

    try:
        raw = await chat_complete_json(
            model=ROUTING_MODEL, system=system, user=user, temperature=0.0
        )
        all_docs = set(docs.keys())
        relevant = [
            line.strip().strip("-").strip()
            for line in raw.strip().split("\n")
            if line.strip() and any(d in line for d in all_docs)
        ]
        return relevant[:4]
    except Exception:
        return []


# ── System prompt ─────────────────────────────────────────────────────────────

SYSTEM_PROMPT = """Sei un esperto analista documentale con accesso completo ai documenti aziendali.

REGOLE ASSOLUTE:
- Rispondi SEMPRE e SOLO in italiano, qualunque sia la lingua del documento
- Usa ESCLUSIVAMENTE le informazioni presenti nei documenti forniti nel contesto
- Non inventare, non stimare, non integrare con conoscenze esterne
- Cita sempre la fonte (nome file, numero slide, numero pagina) per ogni dato riportato
- Se l'informazione non è presente nei documenti, dichiaralo esplicitamente

NUMERI, MISURE E SPECIFICHE TECNICHE — REGOLA CRITICA:
- Riporta i valori ESATTAMENTE come appaiono nel documento: nessuna trasformazione
- Preserva le unità di misura (mm, cm, m, µm, kg, g, ml, l, %, °C, bar, N, pz, €, $, ecc.)
- Non arrotondare mai, non convertire unità, non cambiare il formato
- Tolleranze (±0,5; ±5%; max 3 mm) vanno sempre riportate insieme al valore principale
- Tabelle con misure vanno riportate COMPLETE e FEDELI all'originale
- Se il documento riporta un range (es. 10-15 kg), riporta il range intero

TESTO ESTRATTO DA IMMAGINI ([OCR]):
- Il testo tra [OCR] proviene da immagini, grafici o schemi nel documento
- Trattalo come dato attendibile; segnala la fonte OCR se utile al contesto

STRUTTURA DELLE RISPOSTE:
- Specifiche tecniche → tabella o lista ordinata con unità
- Confronti tra documenti → colonne affiancate con fonte per ogni valore
- Riassunti → struttura gerarchica (sezioni → punti chiave)
- Ricerca di valori specifici → cita il contesto esatto del documento

CAPACITÀ:
- Analisi di presentazioni, relazioni, manuali, capitolati, specifiche tecniche
- Estrazione precisa di dati da tabelle, grafici e immagini (via OCR)
- Confronto tra più documenti sullo stesso argomento
- Ricerca di termini, misure, codici o specifiche esatte
- Sintesi di documenti lunghi mantenendo tutti i dati numerici"""


# ── Build user prompt — usato sia da answer() che dallo streaming ─────────────

async def _build_user_prompt(question: str, context: str) -> str:
    """
    Costruisce il prompt utente per l'answer model, includendo:
      - indice dei documenti indicizzati
      - eventuali annotazioni utente
      - hint dal routing model (documenti probabilmente rilevanti)
      - il contesto di ricerca (chunk + schede)
    """
    loop = asyncio.get_running_loop()
    memory = await loop.run_in_executor(None, load_memory)

    docs = memory.get("documents", {})
    mem_txt = ""

    if docs:
        doc_lines = "\n".join(
            f"  - {name} ({info.get('words', 0):,} parole | {info.get('ext', '')})"
            for name, info in list(docs.items())[:20]
        )
        mem_txt += f"\nDOCUMENTI INDICIZZATI:\n{doc_lines}\n"

    annotations = memory.get("annotations", {})
    if annotations:
        ann_lines = [
            f"  [{fname}] {note}"
            for fname, notes in list(annotations.items())[:5]
            for note in notes[:3]
        ]
        if ann_lines:
            mem_txt += "\nANNOTAZIONI UTENTE:\n" + "\n".join(ann_lines) + "\n"

    relevant = await route_documents(question, memory)
    if relevant:
        mem_txt += "\nDOCUMENTI PROBABILMENTE RILEVANTI PER QUESTA DOMANDA:\n"
        mem_txt += "\n".join(f"  - {d}" for d in relevant) + "\n"

    return f"""{mem_txt}
=== CONTENUTO DOCUMENTI — FONTE PRIMARIA ===
{context}

DOMANDA: {question}

RISPOSTA (obbligatoriamente in italiano — dati solo dalla FONTE PRIMARIA,
numeri e misure ESATTAMENTE come nel documento):"""


# ── Answer (async, non-streaming) ─────────────────────────────────────────────

async def answer(question: str, context: str) -> str:
    """Genera la risposta non-streaming (chiamata da server.py per stream=False)."""
    user_prompt = await _build_user_prompt(question, context)
    return await chat_complete(
        model=ANSWER_MODEL,
        system=SYSTEM_PROMPT,
        user=user_prompt,
        temperature=0.2,
    )


# ── Answer streaming (async iterator) ─────────────────────────────────────────

async def answer_stream(question: str, context: str):
    """
    Genera la risposta in streaming (token-by-token).
    Usato dal server quando stream=True per inoltrare i token via SSE
    direttamente a Open WebUI.
    """
    from llm_client import chat_complete_stream
    user_prompt = await _build_user_prompt(question, context)
    async for chunk in chat_complete_stream(
        model=ANSWER_MODEL,
        system=SYSTEM_PROMPT,
        user=user_prompt,
        temperature=0.2,
    ):
        yield chunk


if __name__ == "__main__":
    folder = sys.argv[1] if len(sys.argv) > 1 else "."
    index_folder(folder)