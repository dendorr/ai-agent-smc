"""Document conversion helpers for the documents agent.

This module converts supported document formats into structured Markdown.
It keeps CPU-bound extraction logic out of ``documents_agent.py`` while
preserving the existing behavior and fallback order.
"""

import re
from pathlib import Path

from config.config import (
    OCR_ENABLED,
    OCR_MIN_TEXT_LEN,
    OCR_PDF_MIN_TEXT_PER_PAGE,
    OCR_PDF_PAGE_RASTER_DPI,
)
from documents.ocr import ocr_image_bytes, ocr_with_glm

# ── Shared helpers ────────────────────────────────────────────────────────────

def _table_to_markdown(table) -> str:
    """
    Convert a Table object (python-pptx or python-docx) to Markdown table.
    Preserves all cell content exactly — no rounding.
    """
    rows = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            rows.append("|" + "|".join(["---"] * len(cells)) + "|")
    return "\n".join(rows)


# ── PPTX → Markdown ───────────────────────────────────────────────────────────

def convert_pptx_to_markdown(filepath) -> str:
    """
    Convert a PPTX to structured Markdown.

    Per slide:
      - Number + title as heading
      - Text boxes (level-aware indentation)
      - Tables → Markdown (exact numbers/units)
      - Embedded images → OCR
      - Speaker notes → blockquote at end of slide
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    p = Path(filepath)
    prs = Presentation(str(filepath))

    lines = [f"# {p.stem}", f"*File: {p.name} — {len(prs.slides)} slide*", ""]

    for slide_num, slide in enumerate(prs.slides, 1):
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except Exception:
            pass

        slide_heading = f"## Slide {slide_num}"
        if title:
            slide_heading += f": {title}"
        lines.append(slide_heading)
        lines.append("")

        for shape in slide.shapes:
            if shape.has_table:
                lines.append(_table_to_markdown(shape.table))
                lines.append("")
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text or text == title:
                        continue
                    level = getattr(para, "level", 0)
                    indent = "  " * level
                    prefix = f"{'#' * (level + 3)} " if level > 0 else ""
                    lines.append(f"{indent}{prefix}{text}")
                lines.append("")
                continue

            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_bytes = shape.image.blob
                    hint = f"[Immagine — slide {slide_num}]"
                    ocr_text = ocr_image_bytes(img_bytes, source_hint=hint)
                    if ocr_text:
                        quoted = ocr_text.replace("\n", "\n> ")
                        lines.append(f"> **{hint}**")
                        lines.append(f"> {quoted}")
                        lines.append("")
            except Exception:
                pass

        try:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip()
            if notes_text:
                quoted = notes_text.replace("\n", "\n> ")
                lines.append("**Note relatore:**")
                lines.append(f"> {quoted}")
                lines.append("")
        except Exception:
            pass

        lines.append("---")
        lines.append("")

    return "\n".join(lines)


# ── DOCX → Markdown ───────────────────────────────────────────────────────────

def convert_docx_to_markdown(filepath) -> str:
    """
    Convert a DOCX to structured Markdown.

    Strategy:
      Tier 1 — markitdown (Microsoft): perfect headings/tables/lists
      Tier 2 — python-docx manual: iterates body in document order
               (paragraphs + tables interleaved), then extracts OCR images

    Numbers, units, tolerances always preserved exactly.
    """
    p = Path(filepath)

    # Tier 1: markitdown
    try:
        from markitdown import MarkItDown
        result = MarkItDown().convert(str(filepath))
        text = result.text_content.strip()
        if len(text) > 100:
            md = f"# {p.stem}\n*File: {p.name}*\n\n{text}"
            img_blocks = _docx_extract_image_ocr(filepath)
            if img_blocks:
                md += "\n\n## Immagini estratte dal documento\n\n" + "\n\n".join(img_blocks)
            return md
    except ImportError:
        pass
    except Exception as e:
        print(f"  [markitdown warning] {p.name}: {e}", flush=True)

    # Tier 2: python-docx manual
    try:
        from docx import Document
        from docx.text.paragraph import Paragraph
        from docx.table import Table as DocxTable

        doc = Document(str(filepath))
        lines = [f"# {p.stem}", f"*File: {p.name}*", ""]

        _HEADING_MAP = {
            "heading 1": "#",
            "heading 2": "##",
            "heading 3": "###",
            "heading 4": "####",
            "heading 5": "#####",
        }

        for child in doc.element.body.iterchildren():
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "p":
                para = Paragraph(child, doc)
                text = para.text.strip()
                style = (para.style.name or "").lower() if para.style else ""

                if not text:
                    lines.append("")
                    continue

                prefix = ""
                for key, md_prefix in _HEADING_MAP.items():
                    if key in style:
                        prefix = md_prefix
                        break

                if prefix:
                    lines.append(f"{prefix} {text}")
                elif "list" in style:
                    level_match = re.search(r"\d", style)
                    indent_level = int(level_match.group()) - 1 if level_match else 0
                    indent = "  " * indent_level
                    lines.append(f"{indent}- {text}")
                else:
                    lines.append(text)

            elif tag == "tbl":
                tbl = DocxTable(child, doc)
                lines.append("")
                lines.append(_table_to_markdown(tbl))
                lines.append("")

        lines.append("")

        img_blocks = _docx_extract_image_ocr(filepath)
        if img_blocks:
            lines.append("## Immagini estratte dal documento")
            lines.extend(img_blocks)

        return "\n".join(lines)

    except Exception as e:
        return f"[DOCX error: {e}]"


def _docx_extract_image_ocr(filepath) -> list:
    """Extract all embedded images from DOCX and run OCR on each."""
    blocks = []
    try:
        from docx import Document
        doc = Document(str(filepath))
        for i, rel in enumerate(doc.part.rels.values(), 1):
            if "image" in rel.reltype:
                try:
                    img_bytes = rel.target_part.blob
                    ocr_text = ocr_image_bytes(
                        img_bytes, source_hint=f"[Immagine {i}]"
                    )
                    if ocr_text and len(ocr_text) > 10:
                        quoted = ocr_text.replace("\n", "\n> ")
                        blocks.append(f"> **Immagine {i}:**\n> {quoted}")
                except Exception:
                    pass
    except Exception:
        pass
    return blocks


# ── PDF → Markdown ────────────────────────────────────────────────────────────

def convert_pdf_to_markdown(filepath) -> str:
    """
    Convert a PDF to structured Markdown.

    Tier 1 — markitdown:  best layout preservation on native PDFs
    Tier 2 — fitz:        text + OCR scanned pages (GLM-OCR)
                          + OCR embedded images + page-by-page tables
    Tier 3 — pdfplumber:  supplementary tables

    Scanned PDFs: pages with < OCR_PDF_MIN_TEXT_PER_PAGE chars are
    rasterized at OCR_PDF_PAGE_RASTER_DPI DPI and sent whole to GLM-OCR.

    All numerical values, units and tolerances preserved exactly.
    """
    import fitz as _fitz

    p = Path(filepath)
    header = f"# {p.stem}\n*File: {p.name}*\n\n"

    # Tier 1: markitdown
    try:
        from markitdown import MarkItDown
        result = MarkItDown().convert(str(filepath))
        text = result.text_content.strip()
        if len(text) > 100:
            return header + text
    except ImportError:
        pass
    except Exception as e:
        print(f"  [markitdown warning] {p.name}: {e}", flush=True)

    # Tier 2: fitz
    fitz_parts = []
    try:
        doc = _fitz.open(str(filepath))
        for page_num, page in enumerate(doc, 1):
            page_lines = [f"## Pagina {page_num}", ""]

            text = page.get_text("text")
            text_clean = text.strip()
            page_was_ocr = False  # True = page handled as scan (GLM-OCR)

            # Page with native text: use it
            if len(text_clean) >= OCR_PDF_MIN_TEXT_PER_PAGE:
                page_lines.append(text_clean)
                page_lines.append("")

            # Page with little text + OCR enabled: probable scan
            elif OCR_ENABLED:
                try:
                    # Full-page rasterization (PDF base = 72 DPI)
                    zoom = OCR_PDF_PAGE_RASTER_DPI / 72
                    mat = _fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat)
                    page_png = pix.tobytes("png")

                    ocr_text = ocr_with_glm(
                        page_png,
                        source_hint=f"[Pagina {page_num} (scansione)]",
                    )
                    if ocr_text and len(ocr_text) > OCR_MIN_TEXT_LEN:
                        page_lines.append(f"[OCR pagina {page_num} — scansione]")
                        page_lines.append(ocr_text)
                        page_lines.append("")
                        page_was_ocr = True
                    elif text_clean:
                        # GLM-OCR failed: use the little native text as fallback
                        page_lines.append(text_clean)
                        page_lines.append("")
                except Exception as e:
                    print(f"  [GLM-OCR warning] pag. {page_num}: {e}", flush=True)
                    if text_clean:
                        page_lines.append(text_clean)
                        page_lines.append("")

            # OCR disabled: use the little native text if present
            elif text_clean:
                page_lines.append(text_clean)
                page_lines.append("")

            # Tables via fitz (PyMuPDF >= 1.23)
            try:
                for tab in page.find_tables().tables:
                    df = tab.to_pandas()
                    md = df.to_markdown(index=False)
                    if md:
                        page_lines.append(f"**Tabella:**\n{md}")
                        page_lines.append("")
            except Exception:
                pass

            # Embedded images → OCR
            # Skip if the page was already handled via GLM-OCR rasterization:
            # the embedded images are already included in the rasterized output —
            # avoids duplicates in ChromaDB.
            if not page_was_ocr:
                for img_info in page.get_images(full=True):
                    try:
                        xref = img_info[0]
                        base_img = doc.extract_image(xref)
                        img_bytes = base_img.get("image", b"")
                        ocr_text = ocr_image_bytes(
                            img_bytes, source_hint=f"[Immagine pag. {page_num}]"
                        )
                        if ocr_text and len(ocr_text) > OCR_MIN_TEXT_LEN:
                            quoted = ocr_text.replace("\n", "\n> ")
                            page_lines.append(f"> **Immagine pag. {page_num}:**\n> {quoted}")
                            page_lines.append("")
                    except Exception:
                        pass

            fitz_parts.append("\n".join(page_lines))
        doc.close()
    except Exception as e:
        fitz_parts.append(f"[fitz error: {e}]")

    # Tier 3: pdfplumber
    plumber_parts = []
    try:
        import pdfplumber
        with pdfplumber.open(str(filepath)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                for table in page.extract_tables() or []:
                    rows = []
                    for i, row in enumerate(table):
                        cells = [str(c).strip() if c else "" for c in row]
                        rows.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            rows.append("|" + "|".join(["---"] * len(cells)) + "|")
                    if rows:
                        plumber_parts.append(
                            f"**Tabella pag. {page_num}:**\n" + "\n".join(rows)
                        )
    except Exception:
        pass

    result = header
    if fitz_parts:
        result += "\n\n---\n\n".join(fitz_parts)
    if plumber_parts:
        result += "\n\n## Tabelle aggiuntive\n\n" + "\n\n".join(plumber_parts)

    return result if len(result) > len(header) + 20 else f"[Impossibile estrarre {p.name}]"
